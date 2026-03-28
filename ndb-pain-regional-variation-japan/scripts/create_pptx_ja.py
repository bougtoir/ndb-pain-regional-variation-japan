#!/usr/bin/env python3
"""Create Japanese PPTX with 1 figure/table per slide.
- Code-generated PNG figures: embedded as images
- Flow diagrams and tables: editable PowerPoint objects
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os

OUTPUT_DIR = '/home/ubuntu/analysis/output/'

# Load regression data
with open(OUTPUT_DIR + 'cpsp_regression_summary.json', 'r') as f:
    reg = json.load(f)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ---- Helper functions ----

def add_title_slide(prs, title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.33), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_image_slide(prs, img_path, title_text, caption_text=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    p.alignment = PP_ALIGN.LEFT
    if os.path.exists(img_path):
        from PIL import Image
        img = Image.open(img_path)
        w, h = img.size
        aspect = w / h
        max_w = Inches(12.5)
        max_h = Inches(5.8)
        if aspect > (12.5 / 5.8):
            pic_w = max_w
            pic_h = int(max_w / aspect)
        else:
            pic_h = max_h
            pic_w = int(max_h * aspect)
        left = int((prs.slide_width - pic_w) / 2)
        top = Inches(1.0)
        slide.shapes.add_picture(img_path, left, top, pic_w, pic_h)
    if caption_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.3), Inches(6.9), Inches(12.73), Inches(0.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = caption_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p2.font.italic = True
    return slide

def add_editable_table_slide(prs, title_text, headers, data_rows, col_widths=None, highlight_rows=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    
    n_rows = len(data_rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        total = 12.0
        col_widths = [total / n_cols] * n_cols
    
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(sum(col_widths))
    height = Inches(min(5.5, 0.45 * n_rows))
    
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for i, cw in enumerate(col_widths):
        table.columns[i].width = Inches(cw)
    
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    
    for r, row_data in enumerate(data_rows):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
            if highlight_rows and r in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xDD, 0xDD)
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    return slide

def add_editable_flow_diagram(prs, title_text, boxes, arrows_between):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    
    shape_refs = {}
    for box in boxes:
        name = box['name']
        left = Inches(box['left'])
        top = Inches(box['top'])
        width = Inches(box['width'])
        height = Inches(box['height'])
        color = box.get('color', (0x1B, 0x3A, 0x5C))
        text = box['text']
        font_size = box.get('font_size', 14)
        
        shape_type = box.get('shape', MSO_SHAPE.ROUNDED_RECTANGLE)
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1)
        
        tf_s = shape.text_frame
        tf_s.word_wrap = True
        tf_s.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p_s = tf_s.paragraphs[0]
            else:
                p_s = tf_s.add_paragraph()
            p_s.text = line
            p_s.font.size = Pt(font_size)
            p_s.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if sum(color) < 400 else RGBColor(0x11, 0x11, 0x11)
            p_s.alignment = PP_ALIGN.CENTER
        
        tf_s.paragraphs[0].space_before = Pt(0)
        shape_refs[name] = {
            'cx': box['left'] + box['width'] / 2,
            'cy': box['top'] + box['height'] / 2,
            'left': box['left'], 'top': box['top'],
            'width': box['width'], 'height': box['height']
        }
    
    for arrow in arrows_between:
        src = shape_refs[arrow['from']]
        dst = shape_refs[arrow['to']]
        direction = arrow.get('direction', 'down')
        
        if direction == 'down':
            x1, y1 = src['cx'], src['top'] + src['height']
            x2, y2 = dst['cx'], dst['top']
        elif direction == 'right':
            x1, y1 = src['left'] + src['width'], src['cy']
            x2, y2 = dst['left'], dst['cy']
        elif direction == 'left':
            x1, y1 = src['left'], src['cy']
            x2, y2 = dst['left'] + dst['width'], dst['cy']
        else:  # up
            x1, y1 = src['cx'], src['top']
            x2, y2 = dst['cx'], dst['top'] + dst['height']
        
        connector = slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        connector.line.width = Pt(2)
    return slide

# ============================================================
# SLIDES
# ============================================================

# --- Slide 1: Title ---
add_title_slide(prs,
    '日本における周術期および慢性疼痛関連処方の地域差',
    'NDBオープンデータを用いた統合的生態学的研究\nPhase 1：急性期疼痛  |  Phase 2：術後遷延性疼痛（CPSP）')

# --- Slide 2: Study Design Flow Diagram (EDITABLE) ---
add_editable_flow_diagram(prs,
    '研究デザイン：データソースと分析フレームワーク',
    boxes=[
        {'name': 'ndb', 'left': 4.5, 'top': 1.0, 'width': 4.5, 'height': 0.8,
         'text': 'NDBオープンデータ 第10回\n（2023年4月〜2024年3月）', 'color': (0x1B, 0x3A, 0x5C), 'font_size': 13},
        
        {'name': 'phase1_data', 'left': 1.0, 'top': 2.3, 'width': 3.5, 'height': 0.7,
         'text': '入院鎮痛薬処方\n（薬効分類114, 811, 821）', 'color': (0x2E, 0x75, 0xB6), 'font_size': 11},
        {'name': 'phase2_data', 'left': 5.0, 'top': 2.3, 'width': 3.5, 'height': 0.7,
         'text': '外来神経障害性疼痛薬\n（PGB, MGB, DLX, TRM, NTP）', 'color': (0xC0, 0x39, 0x2B), 'font_size': 11},
        {'name': 'surgery', 'left': 9.0, 'top': 2.3, 'width': 3.0, 'height': 0.7,
         'text': '入院手術件数\n（Kコード 手術区分）', 'color': (0x27, 0xAE, 0x60), 'font_size': 11},
        
        {'name': 'confounders', 'left': 5.0, 'top': 3.4, 'width': 3.5, 'height': 0.9,
         'text': '交絡疾患プロキシ（外来）：\n糖尿病薬 | 帯状疱疹薬\n抗うつ薬 | 抗不安薬', 'color': (0xF3, 0x9C, 0x12), 'font_size': 10},
        
        {'name': 'phase1', 'left': 1.0, 'top': 4.7, 'width': 3.5, 'height': 0.8,
         'text': 'Phase 1：急性期疼痛指標\n鎮痛薬処方量 / 手術件数', 'color': (0x2E, 0x75, 0xB6), 'font_size': 12},
        {'name': 'phase2', 'left': 5.0, 'top': 4.7, 'width': 3.5, 'height': 0.8,
         'text': 'Phase 2：CPSPプロキシ\n神経障害性疼痛薬 / 手術件数', 'color': (0xC0, 0x39, 0x2B), 'font_size': 12},
        {'name': 'adjusted', 'left': 9.0, 'top': 4.7, 'width': 3.0, 'height': 0.8,
         'text': '調整済CPSP指標\n（交絡除去後の残差）', 'color': (0x8E, 0x44, 0xAD), 'font_size': 12},
        
        {'name': 'integration', 'left': 3.5, 'top': 6.0, 'width': 6.5, 'height': 0.7,
         'text': '統合：急性期 vs 慢性期疼痛の相関 | 東北比較 | 地域ランキング',
         'color': (0x1B, 0x3A, 0x5C), 'font_size': 12},
    ],
    arrows_between=[
        {'from': 'ndb', 'to': 'phase1_data', 'direction': 'down'},
        {'from': 'ndb', 'to': 'phase2_data', 'direction': 'down'},
        {'from': 'ndb', 'to': 'surgery', 'direction': 'down'},
        {'from': 'phase2_data', 'to': 'confounders', 'direction': 'down'},
        {'from': 'phase1_data', 'to': 'phase1', 'direction': 'down'},
        {'from': 'phase2_data', 'to': 'phase2', 'direction': 'down'},
        {'from': 'confounders', 'to': 'adjusted', 'direction': 'down'},
        {'from': 'phase1', 'to': 'integration', 'direction': 'down'},
        {'from': 'phase2', 'to': 'integration', 'direction': 'down'},
        {'from': 'adjusted', 'to': 'integration', 'direction': 'down'},
    ])

# --- Slide 3: Table 1 - Phase 1 results (EDITABLE) ---
add_editable_table_slide(prs,
    '表1. Phase 1：地域ブロック別 入院鎮痛薬/手術指標',
    ['地域ブロック', 'n', '平均（SD）', '範囲', '順位'],
    [
        ['東海', '4', '27.62 (2.18)', '25.20\u201330.07', '1'],
        ['近畿', '6', '30.02 (2.03)', '27.92\u201332.33', '2'],
        ['関東', '7', '33.00 (2.09)', '29.82\u201334.78', '3'],
        ['北陸・甲信越', '6', '35.38 (3.54)', '31.18\u201340.08', '4'],
        ['中国', '5', '35.73 (4.18)', '31.01\u201340.17', '5'],
        ['四国', '4', '36.33 (3.27)', '33.49\u201341.02', '6'],
        ['東北', '6', '39.97 (3.53)', '35.18\u201344.51', '7'],
        ['九州・沖縄', '8', '42.26 (4.33)', '35.82\u201349.75', '8'],
        ['北海道', '1', '46.12 (\u2014)', '46.12', '9'],
        ['全国', '47', '35.78 (5.56)', '25.20\u201349.75', '\u2014'],
    ],
    col_widths=[3.0, 1.0, 2.5, 3.0, 1.5],
    highlight_rows=[6])

# --- Slide 4: Figure 1 - Unadjusted neuropathic bar (PNG) ---
add_image_slide(prs,
    OUTPUT_DIR + 'fig1_neuropathic_unadjusted.png',
    '図1. 都道府県別 外来神経障害性疼痛薬処方量/手術（未調整）',
    '棒グラフ＝神経障害性疼痛薬合計（PGB+MGB+DLX+TRM+NTP）÷入院手術件数。東北（赤枠）が高値側に集中。破線＝全国平均。')

# --- Slide 5: Figure 2 - Confounder correlations (PNG) ---
add_image_slide(prs,
    OUTPUT_DIR + 'fig2_confounder_correlations.png',
    '図2. 神経障害性疼痛薬処方と交絡疾患プロキシの相関',
    '各点＝1都道府県。東北＝赤枠。糖尿病薬が最も強い相関（r=0.87）を示す。')

# --- Slide 6: Figure 3 - Adjusted CPSP index (PNG) ---
add_image_slide(prs,
    OUTPUT_DIR + 'fig3_adjusted_cpsp_index.png',
    '図3. 交絡疾患調整後の都道府県別CPSP指標',
    '残差＝神経障害性疼痛薬処方を糖尿病薬・帯状疱疹薬・抗うつ薬・抗不安薬に回帰した後の残差。東北は調整後に分散。')

# --- Slide 7: Figure 4 - Region unadjusted vs adjusted (PNG) ---
add_image_slide(prs,
    OUTPUT_DIR + 'fig4_region_unadj_vs_adj.png',
    '図4. 地域ブロック別比較：(A) 未調整 vs (B) 交絡疾患調整後',
    '東北（赤枠）は未調整では最高だが、調整後は中位に移動。エラーバー＝標準偏差。')

# --- Slide 8: Table 2 - Regression models (EDITABLE) ---
add_editable_table_slide(prs,
    '表2. 回帰モデル：交絡調整前後の東北効果',
    ['モデル', '従属変数', '東北効果', 'P値', '判定'],
    [
        ['1（未調整）', '神経障害性疼痛薬/手術', f'd = {reg["model1_unadjusted"]["cohens_d"]:.2f}', f'{reg["model1_unadjusted"]["p_value"]:.1e}', '有意 ***'],
        ['2（全交絡調整）', '神経障害性疼痛薬/手術', f'\u03b2 = {reg["model2_adjusted"]["tohoku_coef"]:.1f}', f'{reg["model2_adjusted"]["tohoku_p"]:.3f}', '有意でない'],
        ['3（コア薬）', 'PGB+MGB/手術', f'\u03b2 = {reg["model3_core_neuropathic"]["tohoku_coef"]:.1f}', f'{reg["model3_core_neuropathic"]["tohoku_p"]:.3f}', '有意でない'],
        ['4（神経ブロック）', 'ブロック/手術', f'\u03b2 = {reg["model4_nerve_blocks"]["tohoku_coef"]:.2f}', f'{reg["model4_nerve_blocks"]["tohoku_p"]:.3f}', '有意でない'],
        ['5（統合）', '神経障害性疼痛薬/手術', f'\u03b2 = {reg["model5_integrated"]["tohoku_coef"]:.1f}', f'{reg["model5_integrated"]["tohoku_p"]:.3f}', '有意でない'],
        ['調整済CPSP指標', '残差', f'd = {reg["adjusted_cpsp_test"]["cohens_d"]:.2f}', f'{reg["adjusted_cpsp_test"]["p_value"]:.3f}', '有意でない'],
    ],
    col_widths=[2.5, 2.5, 2.5, 1.5, 2.5],
    highlight_rows=[0])

# --- Slide 9: Figure 5 - Phase 1 vs Phase 2 (PNG) ---
add_image_slide(prs,
    OUTPUT_DIR + 'fig5_phase1_vs_phase2.png',
    '図5. Phase 1（急性期疼痛）vs Phase 2（CPSPプロキシ）の統合',
    '(A) 未調整：r=0.38, P=0.008。(B) 交絡調整後：r=0.29, P=0.052。東北＝赤枠。')

# --- Slide 10: Figure 6 - Model comparison (PNG) ---
add_image_slide(prs,
    OUTPUT_DIR + 'fig6_model_comparison_table.png',
    '図6. 全モデル要約：交絡調整前後の東北効果',
    'モデル1（未調整）は高度有意な過剰を示すが、全調整モデルでは有意でなくなる。')

# --- Slide 11: Supplementary Figure 1 - Heatmap (PNG) ---
add_image_slide(prs,
    OUTPUT_DIR + 'sfig1_heatmap.png',
    '補足図1. 都道府県別 各指標のZ-scoreヒートマップ',
    '各行＝変数、各列＝都道府県（神経障害性疼痛薬/手術で昇順ソート）。赤＝平均以上、青＝平均以下。東北＝赤の垂直線。')

# --- Slide 12: Key Findings Flow (EDITABLE) ---
add_editable_flow_diagram(prs,
    '主要所見の要約',
    boxes=[
        {'name': 'hyp', 'left': 4.0, 'top': 1.0, 'width': 5.5, 'height': 0.7,
         'text': '仮説：東北の「我慢強さ」→ 鎮痛薬使用が少ない',
         'color': (0x95, 0xA5, 0xA6), 'font_size': 14},
        
        {'name': 'p1', 'left': 1.0, 'top': 2.2, 'width': 5.0, 'height': 0.9,
         'text': 'Phase 1：入院周術期鎮痛薬\n東北は高い（d=0.87, P=0.031）\n→ 我慢強さ仮説は棄却',
         'color': (0x2E, 0x75, 0xB6), 'font_size': 12},
        {'name': 'p2u', 'left': 7.0, 'top': 2.2, 'width': 5.5, 'height': 0.9,
         'text': 'Phase 2（未調整）：神経障害性疼痛薬\n東北が最高（d=2.07, P<0.001）\nしかし交絡未調整',
         'color': (0xC0, 0x39, 0x2B), 'font_size': 12},
        
        {'name': 'conf', 'left': 4.0, 'top': 3.6, 'width': 5.5, 'height': 0.7,
         'text': '交絡調整：糖尿病（r=0.87）+ 帯状疱疹 + うつ病 + 不安障害',
         'color': (0xF3, 0x9C, 0x12), 'font_size': 12},
        
        {'name': 'p2a', 'left': 3.0, 'top': 4.8, 'width': 7.5, 'height': 0.9,
         'text': 'Phase 2（調整後）：東北効果は有意でない\n（d=0.44, P=0.323）→ 過剰は糖尿病有病率で説明',
         'color': (0x27, 0xAE, 0x60), 'font_size': 13},
        
        {'name': 'integ', 'left': 2.5, 'top': 6.1, 'width': 8.5, 'height': 0.7,
         'text': '統合：急性期疼痛 ↔ 慢性疼痛に緩やかな正の相関（調整後 r=0.29, P=0.052）',
         'color': (0x8E, 0x44, 0xAD), 'font_size': 13},
    ],
    arrows_between=[
        {'from': 'hyp', 'to': 'p1', 'direction': 'down'},
        {'from': 'hyp', 'to': 'p2u', 'direction': 'down'},
        {'from': 'p2u', 'to': 'conf', 'direction': 'down'},
        {'from': 'conf', 'to': 'p2a', 'direction': 'down'},
        {'from': 'p1', 'to': 'integ', 'direction': 'down'},
        {'from': 'p2a', 'to': 'integ', 'direction': 'down'},
    ])

# --- Slide 13: Confounder Adjustment Concept (EDITABLE) ---
add_editable_flow_diagram(prs,
    '交絡調整フレームワーク（「ノイズ希釈」）',
    boxes=[
        {'name': 'neuro', 'left': 4.5, 'top': 1.0, 'width': 4.5, 'height': 0.7,
         'text': '外来 神経障害性疼痛薬処方\n（プレガバリン、ミロガバリン、デュロキセチン、トラマドール、ノイロトロピン）',
         'color': (0xC0, 0x39, 0x2B), 'font_size': 11},
        
        {'name': 'cpsp', 'left': 0.5, 'top': 2.5, 'width': 3.0, 'height': 0.8,
         'text': 'CPSPシグナル\n（ターゲット）', 'color': (0x27, 0xAE, 0x60), 'font_size': 13},
        {'name': 'dm', 'left': 3.8, 'top': 2.5, 'width': 2.2, 'height': 0.8,
         'text': '糖尿病性\n神経障害', 'color': (0xF3, 0x9C, 0x12), 'font_size': 12},
        {'name': 'hz', 'left': 6.3, 'top': 2.5, 'width': 2.2, 'height': 0.8,
         'text': '帯状疱疹後\n神経痛', 'color': (0xF3, 0x9C, 0x12), 'font_size': 12},
        {'name': 'dep', 'left': 8.8, 'top': 2.5, 'width': 2.0, 'height': 0.8,
         'text': 'うつ病\n（適応外）', 'color': (0xF3, 0x9C, 0x12), 'font_size': 12},
        {'name': 'anx', 'left': 11.1, 'top': 2.5, 'width': 2.0, 'height': 0.8,
         'text': '不安障害\n（適応外）', 'color': (0xF3, 0x9C, 0x12), 'font_size': 12},
        
        {'name': 'dm_proxy', 'left': 3.8, 'top': 4.0, 'width': 2.2, 'height': 0.7,
         'text': '経口血糖降下薬\n（261製剤）', 'color': (0xEB, 0xCB, 0x8B), 'font_size': 10},
        {'name': 'hz_proxy', 'left': 6.3, 'top': 4.0, 'width': 2.2, 'height': 0.7,
         'text': '帯状疱疹抗ウイルス薬\n（47製剤）', 'color': (0xEB, 0xCB, 0x8B), 'font_size': 10},
        {'name': 'dep_proxy', 'left': 8.8, 'top': 4.0, 'width': 2.0, 'height': 0.7,
         'text': '抗うつ薬\n（128製剤）', 'color': (0xEB, 0xCB, 0x8B), 'font_size': 10},
        {'name': 'anx_proxy', 'left': 11.1, 'top': 4.0, 'width': 2.0, 'height': 0.7,
         'text': '抗不安薬\n（112製剤）', 'color': (0xEB, 0xCB, 0x8B), 'font_size': 10},
        
        {'name': 'regression', 'left': 3.0, 'top': 5.3, 'width': 7.5, 'height': 0.7,
         'text': 'OLS回帰：神経障害性疼痛薬 = \u03b2\u2080 + \u03b2\u2081\u00b7糖尿病薬 + \u03b2\u2082\u00b7帯状疱疹薬 + \u03b2\u2083\u00b7抗うつ薬 + \u03b2\u2084\u00b7抗不安薬 + \u03b5',
         'color': (0x1B, 0x3A, 0x5C), 'font_size': 12},
        
        {'name': 'residual', 'left': 4.0, 'top': 6.4, 'width': 5.5, 'height': 0.7,
         'text': '残差（\u03b5）= 調整済CPSP指標\n交絡除去後の「説明されない」神経障害性疼痛薬処方',
         'color': (0x8E, 0x44, 0xAD), 'font_size': 12},
    ],
    arrows_between=[
        {'from': 'cpsp', 'to': 'neuro', 'direction': 'up'},
        {'from': 'dm', 'to': 'neuro', 'direction': 'up'},
        {'from': 'hz', 'to': 'neuro', 'direction': 'up'},
        {'from': 'dep', 'to': 'neuro', 'direction': 'up'},
        {'from': 'anx', 'to': 'neuro', 'direction': 'up'},
        {'from': 'dm', 'to': 'dm_proxy', 'direction': 'down'},
        {'from': 'hz', 'to': 'hz_proxy', 'direction': 'down'},
        {'from': 'dep', 'to': 'dep_proxy', 'direction': 'down'},
        {'from': 'anx', 'to': 'anx_proxy', 'direction': 'down'},
        {'from': 'dm_proxy', 'to': 'regression', 'direction': 'down'},
        {'from': 'hz_proxy', 'to': 'regression', 'direction': 'down'},
        {'from': 'dep_proxy', 'to': 'regression', 'direction': 'down'},
        {'from': 'anx_proxy', 'to': 'regression', 'direction': 'down'},
        {'from': 'regression', 'to': 'residual', 'direction': 'down'},
    ])

# Save
outpath = OUTPUT_DIR + 'figures_JA.pptx'
prs.save(outpath)
print(f'Saved: {outpath}')
