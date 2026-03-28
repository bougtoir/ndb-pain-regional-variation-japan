#!/usr/bin/env python3
"""Fix PPTX z-order and TextBox fill issues.

Problems identified:
1. TextBoxes have no explicit fill → inherit theme fill (often white/opaque),
   blocking shapes/images behind them when clicked.
2. Layout placeholders (Date, Footer, SlideNumber) inherited from Blank layout
   appear as empty selectable areas.
3. On image slides, title/caption TextBoxes should be behind the image.

Fixes:
- Add explicit <a:noFill/> to all TextBox spPr elements
- Remove layout placeholder inheritance where possible
- Reorder spTree children: title TextBox → image/shapes → caption TextBox
  (but since the caption doesn't overlap the image, we keep it on top)
"""

from pptx import Presentation
from lxml import etree
import copy
import sys
import os

NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def qn(tag):
    """Resolve a qualified name like 'a:noFill' to full namespace URI."""
    prefix, local = tag.split(':')
    return f'{{{NSMAP[prefix]}}}{local}'


def set_textbox_nofill(sp_element):
    """Set explicit noFill on a TextBox's spPr so it's transparent."""
    spPr = sp_element.find(qn('a:spPr'))
    if spPr is None:
        # Create spPr if missing
        spPr = etree.SubElement(sp_element, qn('a:spPr'))

    # Remove any existing fill
    for fill_tag in ['a:solidFill', 'a:gradFill', 'a:pattFill', 'a:blipFill', 'a:noFill']:
        existing = spPr.find(qn(fill_tag))
        if existing is not None:
            spPr.remove(existing)

    # Add noFill as first child of spPr (before any xfrm, ln, etc.)
    no_fill = etree.SubElement(spPr, qn('a:noFill'))
    # Move noFill to be after xfrm but before other elements, or first if no xfrm
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(no_fill)
    else:
        spPr.insert(0, no_fill)

    # Also remove any line on the TextBox
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    # Set line to noFill too
    ln_nofill = ln.find(qn('a:noFill'))
    if ln_nofill is None:
        for ln_fill_tag in ['a:solidFill', 'a:gradFill', 'a:pattFill']:
            existing = ln.find(qn(ln_fill_tag))
            if existing is not None:
                ln.remove(existing)
        etree.SubElement(ln, qn('a:noFill'))


def is_textbox(sp_element):
    """Check if an sp element is a TextBox (not an AutoShape)."""
    nvSpPr = sp_element.find(qn('p:nvSpPr'))
    if nvSpPr is None:
        return False
    cNvSpPr = nvSpPr.find(qn('p:cNvSpPr'))
    if cNvSpPr is not None:
        # TextBoxes have txBox="1" attribute
        if cNvSpPr.get('txBox') == '1':
            return True
    return False


def get_shape_name(element):
    """Get the name of a shape element."""
    for path in ['p:nvSpPr/p:cNvPr', 'p:nvPicPr/p:cNvPr', 'p:nvCxnSpPr/p:cNvPr', 'p:nvGrpSpPr/p:cNvPr']:
        parts = path.split('/')
        el = element
        for part in parts:
            el = el.find(qn(part)) if el is not None else None
        if el is not None:
            return el.get('name', 'unknown')
    return 'unknown'


def suppress_layout_placeholders(slide):
    """Add elements to suppress inherited layout placeholders (Date/Footer/SlideNumber)."""
    # In PPTX, to hide inherited placeholders, we can add empty placeholder shapes
    # that override them. But a simpler approach is modifying the slide's XML
    # to explicitly mark them as not shown.
    # Actually, the cleanest way is to set the slide's show master shapes to false,
    # but that would hide ALL master/layout shapes.
    # Instead, let's just make sure the slide layout doesn't have visible placeholders
    # by modifying the layout itself.
    pass  # We'll handle this at the layout level


def reorder_slide_shapes(slide_element):
    """Reorder shapes on a slide so TextBoxes are behind images and AutoShapes.
    
    New z-order (back to front):
    1. nvGrpSpPr (group shape properties - always first, required)
    2. grpSpPr (group shape properties - always second, required)
    3. TextBoxes (title, caption) - BEHIND everything
    4. Pictures (images) - in front of TextBoxes
    5. AutoShapes (rounded rectangles etc.) - in front of pictures
    6. Tables
    7. Connectors (lines/arrows) - MOST FRONT
    """
    sp_tree = slide_element.find(qn('p:cSld') + '/' + qn('p:spTree'))
    if sp_tree is None:
        return

    children = list(sp_tree)
    
    # Separate into categories
    grp_props = []  # nvGrpSpPr, grpSpPr (must stay first)
    textboxes = []
    pictures = []
    autoshapes = []
    tables = []
    connectors = []
    other = []

    for child in children:
        tag = child.tag.split('}')[-1]
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            grp_props.append(child)
        elif tag == 'sp':
            if is_textbox(child):
                textboxes.append(child)
            else:
                autoshapes.append(child)
        elif tag == 'pic':
            pictures.append(child)
        elif tag == 'graphicFrame':
            tables.append(child)
        elif tag == 'cxnSp':
            connectors.append(child)
        else:
            other.append(child)

    # Remove all children
    for child in children:
        sp_tree.remove(child)

    # Re-add in correct z-order (back to front)
    for el in grp_props:
        sp_tree.append(el)
    for el in textboxes:      # TextBoxes at BACK
        sp_tree.append(el)
    for el in pictures:        # Images in middle
        sp_tree.append(el)
    for el in autoshapes:      # AutoShapes in front
        sp_tree.append(el)
    for el in tables:          # Tables in front
        sp_tree.append(el)
    for el in connectors:      # Connectors at FRONT
        sp_tree.append(el)
    for el in other:
        sp_tree.append(el)


def fix_pptx(input_path, output_path):
    """Apply all fixes to a PPTX file."""
    prs = Presentation(input_path)

    # Fix 1: Suppress layout placeholders on the Blank layout
    # Set the Blank layout's placeholders to not show
    blank_layout = prs.slide_layouts[6]  # Blank layout
    layout_element = blank_layout._element
    # Find and remove placeholder shapes from the layout
    layout_sp_tree = layout_element.find(qn('p:cSld') + '/' + qn('p:spTree'))
    if layout_sp_tree is not None:
        for sp in layout_sp_tree.findall(qn('p:sp')):
            nvSpPr = sp.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                nvPr = nvSpPr.find(qn('p:nvPr'))
                if nvPr is not None:
                    ph = nvPr.find(qn('p:ph'))
                    if ph is not None:
                        ph_type = ph.get('type', '')
                        ph_idx = ph.get('idx', '')
                        # Remove Date, Footer, SlideNumber placeholders
                        if ph_type in ('dt', 'ftr', 'sldNum') or ph_idx in ('10', '11', '12'):
                            layout_sp_tree.remove(sp)

    # Fix 2 & 3: For each slide, set noFill on TextBoxes and fix z-order
    for slide_idx, slide in enumerate(prs.slides):
        slide_element = slide._element

        # Find all sp elements that are TextBoxes and set noFill
        sp_tree = slide_element.find(qn('p:cSld') + '/' + qn('p:spTree'))
        if sp_tree is not None:
            for sp in sp_tree.findall(qn('p:sp')):
                if is_textbox(sp):
                    set_textbox_nofill(sp)

        # Reorder shapes: TextBoxes behind, images/shapes in front
        reorder_slide_shapes(slide_element)

    prs.save(output_path)
    print(f'Fixed: {output_path}')


if __name__ == '__main__':
    output_dir = '/home/ubuntu/analysis/output/'
    
    # Fix English version
    fix_pptx(
        output_dir + 'figures_EN.pptx',
        output_dir + 'figures_EN.pptx'
    )
    
    # Fix Japanese version
    fix_pptx(
        output_dir + 'figures_JA.pptx',
        output_dir + 'figures_JA.pptx'
    )
