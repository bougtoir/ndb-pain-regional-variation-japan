#!/usr/bin/env python3
"""Create editable PPTX with English figures for EJP submission.

One figure per slide with title and caption.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'output')
EJP_DIR = os.path.join(OUTPUT_DIR, 'ejp')
os.makedirs(EJP_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

figures = [
    {
        'file': 'fig1_neuropathic_unadjusted_en.png',
        'title': 'Figure 1',
        'caption': 'Outpatient neuropathic pain drug prescribing per surgery by prefecture (unadjusted). '
                   'Tohoku prefectures (red bars) cluster at the high end. Dashed line = national mean.',
    },
    {
        'file': 'fig2_confounder_correlations_en.png',
        'title': 'Figure 2',
        'caption': 'Correlation between outpatient neuropathic pain drug prescribing and confounder disease proxies. '
                   'Diabetes drugs show the strongest correlation (r = 0.87).',
    },
    {
        'file': 'fig3_adjusted_cpsp_index_en.png',
        'title': 'Figure 3',
        'caption': 'Confounder-adjusted CPSP index by prefecture. '
                   'Residuals from regressing neuropathic pain prescribing on four confounder proxies. '
                   'Tohoku prefectures (red borders) are dispersed after adjustment.',
    },
    {
        'file': 'fig4_region_unadj_vs_adj_en.png',
        'title': 'Figure 4',
        'caption': 'Regional comparison of neuropathic pain prescribing: '
                   '(a) unadjusted and (b) after confounder adjustment. '
                   'Tohoku shifts from highest to mid-range after adjustment.',
    },
    {
        'file': 'fig5_phase1_vs_phase2_en.png',
        'title': 'Figure 5',
        'caption': 'Integration of Phase 1 (acute) and Phase 2 (chronic CPSP proxy). '
                   '(a) Unadjusted: r = 0.38, P = 0.008. '
                   '(b) Confounder-adjusted: r = 0.29, P = 0.052.',
    },
    {
        'file': 'sfig1_heatmap_en.png',
        'title': 'Supplementary Figure 1',
        'caption': 'Z-score heatmap of all indices by prefecture. '
                   'Red = above average; blue = below average. '
                   'Tohoku prefectures marked with red vertical lines.',
    },
]

for fig_info in figures:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = fig_info['title']
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT

    # Figure
    fig_path = os.path.join(OUTPUT_DIR, fig_info['file'])
    if os.path.exists(fig_path):
        slide.shapes.add_picture(fig_path, Inches(0.5), Inches(0.8), Inches(12), Inches(5.5))
    else:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'[File not found: {fig_info["file"]}]'
        p2.font.size = Pt(18)
        p2.font.italic = True

    # Caption
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = fig_info['caption']
    p3.font.size = Pt(14)
    p3.font.name = 'Arial'
    p3.alignment = PP_ALIGN.LEFT

outpath = os.path.join(EJP_DIR, 'EJP_figures_EN.pptx')
prs.save(outpath)
print(f'Saved: {outpath}')
