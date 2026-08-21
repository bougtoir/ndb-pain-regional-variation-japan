#!/usr/bin/env python3
"""Create editable PPTX with English figures for JoPPP submission."""
import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'output')
JOPP_DIR = os.path.join(OUTPUT_DIR, 'joppp')
FIGS_DIR = os.path.join(JOPP_DIR, 'figures')
os.makedirs(FIGS_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

figures = [
    {
        'png_file': 'fig1_neuropathic_unadjusted_en.png',
        'tiff_file': 'Fig1_neuropathic_unadjusted.tiff',
        'title': 'Figure 1',
        'caption': 'Outpatient neuropathic pain drug prescribing per surgery by prefecture (unadjusted). '
                   'Tohoku prefectures (red bars) cluster at the high end. '
                   'Dashed line = national mean.',
    },
    {
        'png_file': 'fig2_confounder_correlations_en.png',
        'tiff_file': 'Fig2_confounder_correlations.tiff',
        'title': 'Figure 2',
        'caption': 'Correlation between outpatient neuropathic pain drug prescribing and confounder disease proxies. '
                   'Each dot represents one prefecture. Tohoku prefectures are marked with red borders. '
                   'Diabetes drugs show the strongest correlation (r = 0.87).',
    },
    {
        'png_file': 'fig4_region_unadj_vs_adj_en.png',
        'tiff_file': 'Fig4_region_unadj_vs_adj.tiff',
        'title': 'Figure 3',
        'caption': 'Regional comparison of neuropathic pain prescribing: (a) unadjusted and '
                   '(b) after confounder adjustment. Tohoku (red) shifts from the highest region to '
                   'mid-range after adjustment. Error bars = SD.',
    },
]

for i, fig_info in enumerate(figures, start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = fig_info['title']
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT

    # Figure image
    fig_path = os.path.join(OUTPUT_DIR, fig_info['png_file'])
    if os.path.exists(fig_path):
        slide.shapes.add_picture(fig_path, Inches(0.5), Inches(0.8), Inches(12), Inches(5.5))
    else:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'[File not found: {fig_info["png_file"]}]'
        p2.font.size = Pt(18)
        p2.font.italic = True

    # Caption
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = fig_info['caption']
    p3.font.size = Pt(14)
    p3.font.name = 'Arial'
    p3.alignment = PP_ALIGN.LEFT

    # Copy separate figure files for journal submission
    for src_file, dst_suffix in [(fig_info['png_file'], f'Figure_{i}.png'),
                                 (fig_info['tiff_file'], f'Figure_{i}.tiff')]:
            src = os.path.join(OUTPUT_DIR, src_file)
            if os.path.exists(src):
                shutil.copy2(src, os.path.join(FIGS_DIR, dst_suffix))
                print(f'Copied {src_file} -> {dst_suffix}')

outpath = os.path.join(JOPP_DIR, 'JoPPP_figures_EN.pptx')
prs.save(outpath)
print(f'Saved: {outpath}')
